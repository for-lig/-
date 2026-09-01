import json
import os
import smtplib
import cn
import re

import pymysql
import requests
import tempfile
from datetime import datetime
import pandas as pd
from bs4 import BeautifulSoup
from sklearn.ensemble import IsolationForest
from pyecharts.charts import HeatMap, Line, Grid
from pyecharts import options as opts
from pyecharts.charts import Graph

from email.mime.text import MIMEText
from langchain.tools import tool
from langchain_core.messages import HumanMessage, AIMessage, ToolMessage
from langchain_openai import ChatOpenAI
from zhipuai import ZhipuAI
import gradio as gr
# 热力图 HTML 临时文件路径（用文件中转，彻底绕过作用域问题）
import tempfile

headers = {
    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
}

ZHIPUAI_API_KEY = cn.ZhipuAI  # 替换为你的智谱AI API Key
ZHIPUAI_BASE_URL = "https://open.bigmodel.cn/api/paas/v4/"
WEATHER_API_KEY=cn.weather



@tool
def sq_num(num: int) -> str:
    """求一个数的平方."""
    return f"{num}这个数的平方是{num * num}"

@tool
def baidu(q):
    """
    :param q:
    :return:
    """
    client = ZhipuAI(api_key=cn.ZhipuAI)
    response = client.chat.completions.create(
        model="glm-5.1",
        messages=[
            {
                "role": "system",
                "content": "你是一个有用的AI助手。"
            },
            {
                "role": "user",
                "content": q
            }
        ]
    )
    return response.choices[0].message.content

@tool
def get_weather(location: str) -> str:
    """查询该城市的天气信息."""
    # 调用第三方天气查询工具api
    weather_url = "https://restapi.amap.com/v3/weather/weatherInfo"
    key = WEATHER_API_KEY
    weather_url = f"{weather_url}?key={key}&city={location}"
    text = requests.get(weather_url).text
    # 针对于天气查询的结果做处理，返回给客户
    # print(text)
    info_json = json.loads(text)["lives"][0]
    result = {
        "location": location,
        "date": info_json["reporttime"],
        "weather": info_json["weather"],
        "temperature": f"{info_json['temperature']}°C",
        "humidity": info_json["humidity"]
    }
    # 美观输出

    # 模拟天气 API 调用
    return beautiful_fn(result)


@tool
def send_email(to: str, subject: str, content: str) -> str:
    """向指定邮箱发送邮件.

    Args:
        to: 收件人邮箱地址
        subject: 邮件主题
        content: 邮件正文内容
    """
    mail_host = "smtp.qq.com"  # 设置服务器
    mail_user = cn.USER  # 用户名
    mail_pass = cn.AUTH_CODE  # 口令
    sender = cn.USER

    message = MIMEText(content, 'plain', 'utf-8')
    message['From'] = sender
    message['To'] = to
    message['Subject'] = subject

    try:
        smtpObj = smtplib.SMTP_SSL(mail_host, 465)
        smtpObj.login(mail_user, mail_pass)
        smtpObj.sendmail(sender, [to], message.as_string())
        smtpObj.quit()
        return f"邮件已成功发送至 {to}"
    except smtplib.SMTPException as e:
        return f"邮件发送失败: {e}"

def beautiful_fn(input_txt):
    # 指令
    instruction = """
    任务：查询城市的天气情况
    """
    # 示例（样本）
    example = """
    示例：{'location': '长沙', 'date': '2026-06-09 16:02:05', 'weather': '阴', 'temperature': '25°C', 'humidity': '59'}
    输出：
    城市：长沙
    日期：2026-06-09 16:02:05
    天气：阴
    温度：25°C
    湿度：59
    出行建议：待在家里追剧，如果需要出行，带好雨伞
    """
    # # 输入
    # input_txt = """
    # 晓哥端午回江西吃团圆饭
    # 李四国庆去云南大理洱海找女朋友
    # """
    # 输出
    output_txt = """
    输出
    城市：
    日期：
    天气：
    温度：
    湿度：
    出行建议：
    """

    prompt = f'{instruction}{example}{input_txt}{output_txt}'

    client = ZhipuAI(api_key=cn.ZhipuAI)
    response = client.chat.completions.create(
        model="glm-5.1",
        messages=[
            {
                "role": "system",
                "content": "你是一个有用的AI助手。"
            },
            {
                "role": "user",
                "content": prompt
            }
        ]
    )
    # print(response.choices[0].message.content)
    return response.choices[0].message.content

"""
使用智能体，我们希望通过对话，后台能够自动查询数据，完成信息的整合
对话问题：查询所有人工智能的学生信息 -> 返回所有的学生信息

0.通过大模型，理解对话，决定调用哪个工具（数据处理工具）
1.通过大模型，理解对话问题，生成查询语句   select * from student where major = '人工智能'  / 删除谢芳同学 delete from student where student_name = "谢芳" 
2.执行sql语句，结果返回
"""
# ====================== 数据库操作层 ======================
# 数据库配置
DB_CONFIG = {
    "host": "local",
    "user": "user",
    "password": "123456",
    "port": 3306,
    "database": "user",
    "charset": "utf8mb4"
}


def get_conn():
    """获取数据库连接"""
    return pymysql.connect(**DB_CONFIG)


def query_db(sql, params=None):
    """统一执行SQL查询和非查询操作"""
    conn = get_conn()
    cursor = conn.cursor()

    try:
        cursor.execute(sql, params or ())

        if sql.strip().lower().startswith("select"):
            # 查询操作：返回格式化结果
            rows = cursor.fetchall()
            cols = [desc[0] for desc in cursor.description]
            res = [" | ".join(cols), "-" * 50]
            for row in rows:
                res.append(" | ".join(str(x) for x in row))
            return "\n".join(res)
        else:
            # 非查询操作：提交事务并返回影响行数
            conn.commit()
            return f"操作成功，影响 {cursor.rowcount} 行"
    finally:
        cursor.close()
        conn.close()


# 2. 定义唯一的SQL查询工具
@tool
def execute_sql_query(query: str) -> str:
    """
    执行学生和成绩相关的数据库操作
    支持查询、添加、修改、删除学生信息和成绩
    参数：query - 用户的自然语言问题
    """
    print(f"[工具调用] 执行SQL查询: {query}")

    # SQL生成系统提示词
    sql_prompt = f"""
    你是一个专业的MySQL 8.0 SQL生成器，**只返回合法的SQL语句**。

    表结构：
    students: student_id, student_name, student_no, gender, major, class_name, email, phone
    scores: score_id, student_id, course_name, score, semester, exam_time

    规则：
    1. 查询成绩必须使用JOIN关联students和scores表
    2. 优先使用student_name作为查询条件
    3. 支持AVG、COUNT、SUM等聚合函数
    4. 支持INSERT、UPDATE、DELETE、SELECT所有操作
    5. 绝对不返回任何自然语言解释、注释或说明
    6. 禁止返回中文内容 
    """

    # 调用大模型生成SQL
    sql_llm = ChatOpenAI(
        model="glm-4-plus",
        api_key=ZHIPUAI_API_KEY,
        base_url=ZHIPUAI_BASE_URL,
        temperature=0.1,
        timeout=30
    )
    response = sql_llm.invoke([
        ("system", sql_prompt),
        ("user", query)
    ])

    # 提取并清理SQL（多重防护，避免自然语言混入）
    sql = response.content.strip()
    sql = sql.replace("```sql", "").replace("```", "").strip()

    """
    会出现的场景：
    \-\- 这是一段注释
    SELECT id,name FROM user WHERE age>18;

    随便一段文字 abc123 UPDATE student SET score=90
    """
    # 正则提取纯SQL（终极防护）
    # re.I（re.IGNORECASE）：忽略大小写，select / Select / SELECT 全都能匹配
    # re.DOTALL：让 . 可以匹配换行符，支持跨多行 SQL 语句
    sql_match = re.search(r"(SELECT|INSERT|UPDATE|DELETE).*", sql, re.I | re.DOTALL)
    if sql_match:
        sql = sql_match.group(0).strip()

    # 最终校验
    if not sql or not sql.upper().startswith(("SELECT", "INSERT", "UPDATE", "DELETE")):
        return "错误：无法生成合法的SQL语句，请重新描述您的问题。"

    print(f"[生成SQL] {sql}")
    return query_db(sql)

def _get_db_conn():
    """获取 MySQL 连接"""
    return pymysql.connect(**DB_CONFIG)

def _init_stock_table():
    """初始化股票数据表"""
    conn = _get_db_conn()
    try:
        with conn.cursor() as cursor:
            cursor.execute("""
                CREATE TABLE IF NOT EXISTS stock_industry_flow (
                    id INT AUTO_INCREMENT PRIMARY KEY,
                    fetch_time DATETIME NOT NULL,
                    industry VARCHAR(64) NOT NULL,
                    industry_index VARCHAR(32),
                    pct_change VARCHAR(16),
                    funds_in DECIMAL(12,2),
                    funds_out DECIMAL(12,2),
                    net_amount DECIMAL(12,2),
                    num_companies INT,
                    leading_stock VARCHAR(64),
                    leading_pct VARCHAR(16),
                    current_price VARCHAR(16)
                ) ENGINE=InnoDB DEFAULT CHARSET=utf8mb4;
            """)
        conn.commit()
    finally:
        conn.close()

def _init_question_bank_table():
    """初始化搜题库表"""
    conn = _get_db_conn()
    try:
        with conn.cursor() as cursor:
            cursor.execute("""
                CREATE TABLE IF NOT EXISTS question_bank (
                    id INT AUTO_INCREMENT PRIMARY KEY,
                    create_time DATETIME NOT NULL,
                    subject VARCHAR(64) NOT NULL,
                    question TEXT NOT NULL,
                    answer TEXT NOT NULL
                ) ENGINE=InnoDB DEFAULT CHARSET=utf8mb4;
            """)
        conn.commit()
    finally:
        conn.close()

@tool
def search_book(book_keyword: str) -> str:
    """
    无数据库图书查询，根据关键词查询图书信息（书名、作者、简介、分类、推荐指数）。
    :param book_keyword: 书名、作者、图书分类关键词
    """
    client = ZhipuAI(api_key=cn.ZhipuAI)
    prompt = f"""
    你是图书馆图书检索助手，根据关键词查询图书，返回清晰结构化图书信息，包含书名、作者、简介、分类、推荐指数。
    用户查询关键词：{book_keyword}
    如果用户需要新增/修改图书，告知当前无数据库，仅能查询图书科普信息。
    """
    resp = client.chat.completions.create(model="glm-4-flash", messages=[{"role": "user", "content": prompt}])
    return resp.choices[0].message.content


@tool
def analyze_file(file_path: str) -> str:
    """
    分析本地文件内容，支持txt、csv、json、md、py等文本格式文件。
    :param file_path: 文件的绝对路径或相对路径
    """
    try:
        if not os.path.exists(file_path):
            return f"错误：文件不存在 - {file_path}"

        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        file_info = {"文件名": os.path.basename(file_path), "文件路径": file_path, "文件大小": f"{os.path.getsize(file_path)} bytes", "文件类型": ext}

        if ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.java', '.c', '.cpp']:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            content_preview = content[:5000] + "\n...（内容已截断）" if len(content) > 5000 else content
            analysis_prompt = f"请分析以下文件内容并提供摘要：\n文件信息：{file_info}\n文件内容：\n{content_preview}\n请提供：1.文件内容摘要 2.主要主题或关键点 3.如果是代码文件，说明其功能"
            client = ZhipuAI(api_key=cn.ZhipuAI)
            resp = client.chat.completions.create(model="glm-4-flash", messages=[{"role": "system", "content": "你是一个专业的文件分析助手。"}, {"role": "user", "content": analysis_prompt}])
            return f"=== 文件分析报告 ===\n\n{resp.choices[0].message.content}\n\n=== 文件基本信息 ===\n" + "\n".join([f"{k}: {v}" for k, v in file_info.items()])

        elif ext == '.csv':
            import csv
            with open(file_path, 'r', encoding='utf-8') as f:
                rows = list(csv.reader(f))
            preview = "\n".join([",".join(row) for row in rows[:10]])
            if len(rows) > 10:
                preview += f"\n...（共{len(rows)}行，仅显示前10行）"
            client = ZhipuAI(api_key=cn.ZhipuAI)
            resp = client.chat.completions.create(model="glm-4-flash", messages=[{"role": "system", "content": "你是一个专业的数据分析助手。"}, {"role": "user", "content": f"分析CSV文件：\n文件名：{file_info['文件名']}\n行数：{len(rows)}\n数据预览：\n{preview}\n请提供数据结构说明、主要字段含义、数据特点"}])
            return f"=== CSV分析报告 ===\n\n{resp.choices[0].message.content}\n\n=== 基本信息 ===\n" + "\n".join([f"{k}: {v}" for k, v in file_info.items()])

        elif ext == '.json':
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            json_str = json.dumps(data, ensure_ascii=False, indent=2)
            json_preview = json_str[:5000] + "\n...（截断）" if len(json_str) > 5000 else json_str
            client = ZhipuAI(api_key=cn.ZhipuAI)
            resp = client.chat.completions.create(model="glm-4-flash", messages=[{"role": "system", "content": "你是JSON数据分析助手。"}, {"role": "user", "content": f"分析JSON文件：\n{json_preview}\n请提供结构说明、主要字段含义、数据用途推测"}])
            return f"=== JSON分析报告 ===\n\n{resp.choices[0].message.content}\n\n=== 基本信息 ===\n" + "\n".join([f"{k}: {v}" for k, v in file_info.items()])

        else:
            return f"不支持的文件类型：{ext}\n支持：.txt, .md, .py, .js, .html, .css, .java, .c, .cpp, .csv, .json"
    except Exception as e:
        return f"文件分析失败：{str(e)}"

@tool
def translate_text(text: str, target_lang: str = "中文") -> str:
    """
    多语言翻译工具，支持中、英、日、韩、法、德、西等多种语言
    :param text: 需要翻译的文本
    :param target_lang: 目标语言，默认为中文
    """
    lang_map = {
        "中文": "Chinese", "英文": "English", "英语": "English",
        "日文": "Japanese", "日语": "Japanese", "韩文": "Korean", "韩语": "Korean",
        "法文": "French", "法语": "French", "德文": "German", "德语": "German",
        "西班牙文": "Spanish", "西班牙语": "Spanish", "俄文": "Russian", "俄语": "Russian"
    }
    target = lang_map.get(target_lang, target_lang)

    client = ZhipuAI(api_key=cn.ZhipuAI)
    prompt = f"请将以下文本翻译成{target}，只返回翻译结果，不要添加任何解释：\n\n原文：{text}\n\n翻译结果："

    try:
        resp = client.chat.completions.create(
            model="glm-4-flash",
            messages=[
                {"role": "system", "content": "你是一个专业的翻译助手，只返回翻译结果。"},
                {"role": "user", "content": prompt}
            ]
        )
        translated = resp.choices[0].message.content.strip()
        return f"**翻译结果**\n---\n原文：{text}\n目标语言：{target_lang}\n译文：{translated}\n---"
    except Exception as e:
        return f"翻译失败：{str(e)}"

def _save_records_to_db(records):
    """写入股票数据到MySQL"""
    _init_stock_table()
    fetch_time = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    conn = _get_db_conn()
    try:
        with conn.cursor() as cursor:
            for r in records:
                def safe_float(val):
                    try:
                        return float(str(val).replace('%', '').replace('+', '').strip())
                    except:
                        return None

                cursor.execute("""
                    INSERT INTO stock_industry_flow
                    (fetch_time, industry, industry_index, pct_change, funds_in, funds_out, net_amount, num_companies, leading_stock, leading_pct, current_price)
                    VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s)
                """, (fetch_time, r.get('Industry', ''), r.get('Industry_Index', ''), r.get('Pct Change', ''),
                      safe_float(r.get('Funds In (B)')), safe_float(r.get('Funds Out (B)')),
                      safe_float(r.get('Net Amount (B)')),
                      safe_float(r.get('No. of Companies')), r.get('Leading Stock', ''),
                      r.get('Pct Change (Leading)', ''), r.get('Current Price (¥)', '')))
        conn.commit()
    finally:
        conn.close()


def _load_latest_records_from_db():
    """从MySQL读取最新股票数据"""
    conn = _get_db_conn()
    try:
        with conn.cursor(pymysql.cursors.DictCursor) as cursor:
            cursor.execute("""
                SELECT * FROM stock_industry_flow
                WHERE fetch_time = (SELECT MAX(fetch_time) FROM stock_industry_flow)
                ORDER BY id ASC
            """)
            return cursor.fetchall()
    finally:
        conn.close()


def fn_stock(input_txt):
    """AI解读股票数据"""
    prompt = f"""
    任务：解读A股行业资金流向数据，给出摘要和投资参考。
    全量数据：{input_txt}
    字段说明：Industry=行业, Pct Change=涨跌幅, Funds In=流入(亿), Funds Out=流出(亿), Net Amount=净额(亿), Leading Stock=龙头股
    输出：净流入前五行业、净流出前五行业、市场整体资金趋势、涨幅最大龙头股、市场情绪判断、一句话建议
    """
    client = ZhipuAI(api_key=cn.ZhipuAI)
    resp = client.chat.completions.create(
        model='glm-4-plus',
        messages=[{'role': 'system', 'content': '你是专业的股票市场分析师'}, {'role': 'user', 'content': prompt}]
    )
    return resp.choices[0].message.content


@tool
def get_stock_info() -> str:
    """获取A股市场实时行业资金流向数据"""
    url = 'https://data.10jqka.com.cn/funds/gnzjl/'
    try:
        response = requests.get(url, headers=headers, timeout=10)
    except Exception as e:
        return f"股票数据抓取失败：{e}"

    soup = BeautifulSoup(response.content, 'html.parser')
    table = soup.find('table', attrs={'class': 'm-table J-ajax-table'})
    if not table:
        return "未能找到股票数据表格。"

    col_names = ['No.', 'Industry', 'Industry_Index', 'Pct Change', 'Funds In (B)', 'Funds Out (B)',
                 'Net Amount (B)', 'No. of Companies', 'Leading Stock', 'Pct Change (Leading)', 'Current Price (¥)']
    records = []
    for row in table.find("tbody").find_all("tr"):
        values = [cell.text.strip().split('\n')[0] for cell in row.find_all("td")]
        if len(values) == len(col_names):
            records.append(dict(zip(col_names, values)))

    if not records:
        return "解析到的股票数据为空。"

    _save_records_to_db(records)
    return fn_stock(records)


@tool
def detect_stock_anomaly() -> str:
    """检测A股行业资金流向中的异常行业"""
    rows = _load_latest_records_from_db()
    if not rows:
        return "数据库中暂无股票数据，请先查询一次股票行情。"

    records = [{'Industry': r['industry'], 'Pct Change': r['pct_change'], 'Funds In (B)': str(r['funds_in'] or 0),
                'Funds Out (B)': str(r['funds_out'] or 0), 'Net Amount (B)': str(r['net_amount'] or 0),
                'Industry_Index': str(r['industry_index'] or 0), 'No. of Companies': str(r['num_companies'] or 0),
                'Leading Stock': r['leading_stock']} for r in rows]

    df = pd.DataFrame(records)
    df['Pct Change Num'] = df['Pct Change'].str.rstrip('%').astype(float)
    for col in ['Funds In (B)', 'Funds Out (B)', 'Net Amount (B)', 'Industry_Index', 'No. of Companies']:
        df[col] = pd.to_numeric(df[col], errors='coerce').fillna(0)

    features = df[
        ['Industry_Index', 'Pct Change Num', 'Funds In (B)', 'Funds Out (B)', 'Net Amount (B)', 'No. of Companies']]
    df['anomaly'] = IsolationForest(contamination=0.05, random_state=42).fit_predict(features)

    anomaly_list = df[df['anomaly'] == -1][['Industry', 'Pct Change', 'Net Amount (B)', 'Leading Stock']].to_dict(
        orient='records')
    normal_top3 = df[df['anomaly'] == 1].head(3)[['Industry', 'Pct Change', 'Net Amount (B)', 'Leading Stock']].to_dict(
        orient='records')

    prompt = f"解读A股异常检测结果。异常行业：{anomaly_list}。正常参考：{normal_top3}。输出：异常行业列表、原因分析、是否值得关注、风险提示。"
    client = ZhipuAI(api_key=cn.ZhipuAI)
    resp = client.chat.completions.create(model='glm-4-plus',
                                          messages=[{'role': 'system', 'content': '你是专业的股票市场风险分析师'},
                                                    {'role': 'user', 'content': prompt}])
    return resp.choices[0].message.content


@tool
def show_stock_heatmap() -> str:
    """生成A股行业资金净额热力图"""
    rows = _load_latest_records_from_db()
    if not rows:
        return "HEATMAP_ERROR|数据库中暂无股票数据，请先查询一次股票行情。"

    records = [{'Industry': r['industry'], 'Net Amount (B)': float(r['net_amount'] or 0), 'Pct Change': r['pct_change']}
               for r in rows]
    df = pd.DataFrame(records)

    industries = df['Industry'].tolist()
    net_amounts = df['Net Amount (B)'].tolist()

    heatmap = HeatMap().add_xaxis(industries).add_yaxis("", ["净额 (亿)"],
                                                        [[i, 0, net_amounts[i]] for i in range(len(df))],
                                                        label_opts=opts.LabelOpts(is_show=False)).set_global_opts(
        title_opts=opts.TitleOpts(title="行业资金净额热力图"),
        visualmap_opts=opts.VisualMapOpts(min_=min(net_amounts), max_=max(net_amounts), orient="horizontal",
                                          pos_top="5%"),
        xaxis_opts=opts.AxisOpts(axislabel_opts={"rotate": 45, "interval": 0}))

    line = Line().add_xaxis(industries).add_yaxis("净额 (亿)", net_amounts, is_smooth=True,
                                                  linestyle_opts=opts.LineStyleOpts(width=3,
                                                                                    color="#4a4a4a")).set_global_opts(
        xaxis_opts=opts.AxisOpts(axislabel_opts={"rotate": 45, "interval": 0}))

    grid = Grid(init_opts=opts.InitOpts(width="100%", height="600px"))
    grid.add(heatmap, grid_opts=opts.GridOpts(pos_left="5%", pos_right="5%", pos_top="15%", pos_bottom="25%"))
    grid.add(line, grid_opts=opts.GridOpts(pos_left="5%", pos_right="5%", pos_top="15%", pos_bottom="25%"))

    tmp_file = os.path.join(tempfile.gettempdir(), "stock_heatmap.html")
    grid.render(tmp_file)
    return "HEATMAP_READY"


@tool
def generate_knowledge_graph(file_path: str) -> str:
    """分析文档并生成知识图谱"""
    if not os.path.exists(file_path):
        return f"KGRAPH_ERROR|文件不存在：{file_path}"

    ext = os.path.splitext(file_path)[1].lower()
    content = ""
    try:
        if ext == '.pdf':
            import PyPDF2
            with open(file_path, 'rb') as f:
                for page in PyPDF2.PdfReader(f).pages:
                    content += (page.extract_text() or "") + "\n"
        elif ext in ['.pptx']:
            from pptx import Presentation
            for slide in Presentation(file_path).slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): content += shape.text + "\n"
        else:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
    except Exception as e:
        return f"KGRAPH_ERROR|文件读取失败: {e}"

    if not content.strip():
        return "KGRAPH_ERROR|未能提取到有效文本。"

    prompt = f"""分析以下文本，提取核心概念及关系，构建知识图谱。
    输出要求：必须是合法纯JSON，不要markdown符号。格式：
    {{"nodes": [{{"name": "中心主题", "symbolSize": 60}}], "links": [{{"source": "中心主题", "target": "子主题", "value": "关联"}}]}}
    提取10-20个核心节点。

    文本：{content[:4000]}"""

    try:
        client = ZhipuAI(api_key=cn.ZhipuAI)
        resp = client.chat.completions.create(model='glm-4-plus', messages=[{'role': 'user', 'content': prompt}],
                                              temperature=0.1)
        result_text = re.sub(r"```json|```", "", resp.choices[0].message.content.strip()).strip()
        graph_data = json.loads(result_text)
    except Exception as e:
        return f"KGRAPH_ERROR|JSON解析失败：{str(e)}"

    nodes = graph_data.get("nodes", [])
    links = graph_data.get("links", [])

    kg = Graph(init_opts=opts.InitOpts(width="100%", height="650px", theme="light")).add(
        "知识网络", nodes, links, layout="force", repulsion=6000, gravity=0.1,
        is_roam=True, is_focusnode=True, is_draggable=True,
        edge_label=opts.LabelOpts(is_show=True, position="middle", formatter="{c}", color="#888", font_size=12),
        label_opts=opts.LabelOpts(is_show=True, font_size=14, font_weight="bold"),
        tooltip_opts=opts.TooltipOpts(is_show=True, formatter="{b}")
    ).set_global_opts(title_opts=opts.TitleOpts(title="📄 文档核心知识图谱", pos_left="center"))

    tmp_file = os.path.join(tempfile.gettempdir(), "kgraph_cache.html")
    kg.render(tmp_file)
    return f"KGRAPH_READY|{tmp_file}"


@tool
def save_to_question_bank(subject: str, question: str, answer: str) -> str:
    """将问答记录存入搜题库"""
    _init_question_bank_table()
    conn = _get_db_conn()
    try:
        with conn.cursor() as cursor:
            cursor.execute(
                "INSERT INTO question_bank (create_time, subject, question, answer) VALUES (%s, %s, %s, %s)",
                (datetime.now().strftime('%Y-%m-%d %H:%M:%S'), subject, question, answer)
            )
        conn.commit()
        return f"✅ **已自动存入专属搜题库（科目：{subject}）**\n\n---\n\n{answer}"
    except Exception as e:
        return f"❌ 存入失败：{e}\n\n---\n\n**详细解答：**\n{answer}"
    finally:
        conn.close()


@tool
def summarize_question_bank() -> str:
    """总结搜题库，生成学习报告"""
    _init_question_bank_table()
    conn = _get_db_conn()
    try:
        with conn.cursor(pymysql.cursors.DictCursor) as cursor:
            cursor.execute("SELECT subject, question, answer FROM question_bank ORDER BY subject")
            rows = cursor.fetchall()
    finally:
        conn.close()

    if not rows:
        return "你的搜题库目前是空的。"

    db_content = "".join([f"【科目】{r['subject']}\n【问题】{r['question']}\n【回答】{r['answer']}\n---\n" for r in rows])

    prompt = f"你是辅导老师。根据学生搜题库记录生成结构化学习总结。按科目分类、提炼薄弱点、给出学习建议。\n\n记录：{db_content[:6000]}"
    client = ZhipuAI(api_key=cn.ZhipuAI)
    resp = client.chat.completions.create(model='glm-4-plus',
                                          messages=[{'role': 'system', 'content': '你是专业的学习总结助手'},
                                                    {'role': 'user', 'content': prompt}], temperature=0.3)
    summary = resp.choices[0].message.content

    filepath = os.path.join(tempfile.gettempdir(), "专属学习总结报告.md")
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(summary)
    return f"SUMMARY_READY|{filepath}"

@tool
def get_music_preview(song_name: str) -> str:
    """搜索歌曲并返回30秒试听片段链接。"""
    search_url = "https://itunes.apple.com/search"
    resp = requests.get(search_url, params={"term": song_name, "media": "music", "limit": 1}, timeout=10)
    if resp.status_code != 200:
        return f"搜索失败：{resp.status_code}"
    results = resp.json().get("results", [])
    if not results:
        return f"未找到歌曲：{song_name}"
    track = results[0]
    name = track.get("trackName", "未知歌曲")
    artist = track.get("artistName", "未知歌手")
    preview_url = track.get("previewUrl")
    if not preview_url:
        return f"找到歌曲《{name}》- {artist}，但该歌曲暂无试听片段。"
    
    # 下载音频到本地临时文件，使用安全的文件名
    try:
        audio_resp = requests.get(preview_url, timeout=10)
        audio_resp.raise_for_status()
        # 使用 trackId 作为文件名，避免特殊字符问题
        track_id = track.get("trackId", str(hash(name + artist)))
        tmp_file = os.path.join(tempfile.gettempdir(), f"audio_{track_id}.m4a")
        with open(tmp_file, 'wb') as f:
            f.write(audio_resp.content)
        return f"MUSIC_PREVIEW|{name}|{artist}|{tmp_file}"
    except Exception as e:
        return f"下载音频失败：{str(e)}"



# ==================== 多任务 + 流式输出 ====================

tools_list = [sq_num, baidu, get_weather, send_email, execute_sql_query, search_book, analyze_file, translate_text,
              get_stock_info, detect_stock_anomaly, show_stock_heatmap,
              generate_knowledge_graph, save_to_question_bank, summarize_question_bank, get_music_preview]

tool_map = {t.name: t for t in tools_list}



def get_name(question):
    try:
        client = ZhipuAI(api_key=cn.ZhipuAI)
        resp = client.chat.completions.create(
            model="glm-4-flash",
            messages=[
                {"role": "system", "content": "用不超过8个字概括用户意图，只输出文字，不加符号"},
                {"role": "user", "content": question}
            ]
        )
        name = resp.choices[0].message.content.strip().strip('"').strip("'")
        return name if name else "新对话"
    except Exception:
        return "新对话"


def update_radio(tasks_dict, current_name):
    """统一的Radio更新：返回gr.update."""
    keys = list(tasks_dict.keys())
    if current_name not in keys:
        current_name = keys[-1] if keys else "新对话"
    return gr.update(choices=keys, value=current_name)


def chat_stream(prompt, chatbot, tasks_dict, current_name, file_upload=None, task_visuals=None):
    """核心：多任务记忆 + 逐字流式输出 + 特殊标记解析"""
    if task_visuals is None:
        task_visuals = {}
    if not current_name or current_name not in tasks_dict:
        current_name = list(tasks_dict.keys())[-1] if tasks_dict else "新对话"

    history = tasks_dict[current_name]
    history.append({"role": "user", "content": prompt})

    # 第一条消息 → 自动命名
    if len(history) == 1:
        new_title = get_name(prompt)
        old_name = current_name
        tasks_dict[new_title] = tasks_dict.pop(old_name)
        # 同步更新 task_visuals 的键
        if old_name in task_visuals:
            task_visuals[new_title] = task_visuals.pop(old_name)
        current_name = new_title

    # 构建消息（记忆）
    lc_msgs = []
    for m in history:
        c = m.get("content", "")
        if not c:
            continue
        if m["role"] == "user":
            lc_msgs.append(HumanMessage(content=c))
        else:
            lc_msgs.append(AIMessage(content=c))

    # 系统提示词（特殊规则）
    system_msg = """你有以下特殊输出规则：
【规则1-热力图】调用show_stock_heatmap后返回HEATMAP_READY，必须原样输出。
【规则2-知识图谱】调用generate_knowledge_graph后返回KGRAPH_READY|路径，必须原样输出。
【规则3-搜题库】调用save_to_question_bank后，回复中必须包含详细解答。
【规则4-总结报告】调用summarize_question_bank后返回SUMMARY_READY|路径，必须原样输出。
【规则5-音乐试听】调用get_music_preview后返回MUSIC_PREVIEW|name|artist|url，必须原样输出。"""

    # 如果上传了文件，在提示中附加文件路径
    actual_prompt = prompt
    if file_upload:
        actual_prompt += f"\n\n（用户上传了文件：{file_upload}，请先调用analyze_file总结，再调用generate_knowledge_graph生成知识图谱。）"

    # 重新构建带系统提示的消息
    lc_msgs = [HumanMessage(content=system_msg)] + lc_msgs[:-1] + [HumanMessage(content=actual_prompt)]

    model = ChatOpenAI(
        model="glm-4-plus", api_key=ZHIPUAI_API_KEY,
        base_url=ZHIPUAI_BASE_URL, temperature=0.1, timeout=30
    )
    model_bind = model.bind_tools(tools_list)

    full_reply = ""
    history.append({"role": "assistant", "content": ""})

    # 存储工具返回的原始标记（避免被模型重写）
    tool_raw_outputs = []

    for _ in range(5):
        text_part = ""
        tool_calls = {}

        for chunk in model_bind.stream(lc_msgs):
            # 处理文本内容
            if chunk.content:
                text_part += chunk.content
                full_reply += chunk.content
                history[-1]["content"] = full_reply
                yield "", list(history), update_radio(tasks_dict, current_name), current_name, None, None, None

            # 积累工具调用信息
            if hasattr(chunk, 'tool_call_chunks') and chunk.tool_call_chunks:
                for tc_chunk in chunk.tool_call_chunks:
                    idx = tc_chunk.get('index', 0)
                    if idx not in tool_calls:
                        tool_calls[idx] = {'id': '', 'name': '', 'args': ''}
                    if tc_chunk.get('id'):
                        tool_calls[idx]['id'] = tc_chunk['id']
                    if tc_chunk.get('name'):
                        tool_calls[idx]['name'] = tc_chunk['name']
                    if tc_chunk.get('args'):
                        tool_calls[idx]['args'] += tc_chunk['args']

        ai_msg = AIMessage(content=text_part)
        if tool_calls:
            ai_msg.tool_calls = []
            for i in sorted(tool_calls):
                tc = tool_calls[i]
                try:
                    args = json.loads(tc["args"]) if tc["args"] else {}
                except Exception:
                    args = {}
                ai_msg.tool_calls.append({"id": tc["id"], "name": tc["name"], "args": args})
        lc_msgs.append(ai_msg)

        if not tool_calls:
            break

        for tc in ai_msg.tool_calls:
            fn = tool_map.get(tc["name"])
            if fn:
                try:
                    r = fn.invoke(tc["args"])
                    # 保存工具原始返回值（包含特殊标记）
                    tool_raw_outputs.append(str(r))
                    lc_msgs.append(ToolMessage(content=str(r), tool_call_id=tc["id"]))
                except Exception as e:
                    lc_msgs.append(ToolMessage(content=f"失败: {e}", tool_call_id=tc["id"]))

        full_reply = ""
        history[-1]["content"] = ""

    # ===== 标记解析 =====
    preview_url = None
    heatmap_html = None
    download_path = None

    # 合并所有工具原始输出，优先从这里解析标记
    tool_output_str = "\n".join(tool_raw_outputs)

    # 1. 解析热力图
    if "HEATMAP_READY" in tool_output_str or "HEATMAP_READY" in full_reply:
        source_str = tool_output_str if "HEATMAP_READY" in tool_output_str else full_reply
        tmp_file = os.path.join(tempfile.gettempdir(), "stock_heatmap.html")
        if os.path.exists(tmp_file):
            with open(tmp_file, 'r', encoding='utf-8') as f:
                raw = f.read()
            os.remove(tmp_file)
            import html
            heatmap_html = f'<iframe srcdoc="{html.escape(raw)}" width="100%" height="650px" style="border:none;"></iframe>'
        full_reply = full_reply.replace("HEATMAP_READY", "").strip()
        if not full_reply:
            full_reply = "📊 已生成实时股票行业资金净额热力图。"

    # 2. 解析知识图谱
    if "KGRAPH_READY|" in tool_output_str or "KGRAPH_READY|" in full_reply:
        try:
            source_str = tool_output_str if "KGRAPH_READY|" in tool_output_str else full_reply
            marker = source_str.index("KGRAPH_READY|")
            file_path = source_str[marker + 13:].strip().split('\n')[0].strip()
            if os.path.exists(file_path):
                with open(file_path, 'r', encoding='utf-8') as f:
                    raw = f.read()
                os.remove(file_path)
                import html
                heatmap_html = f'<iframe srcdoc="{html.escape(raw)}" width="100%" height="650px" style="border:none;"></iframe>'
            full_reply = full_reply.replace("KGRAPH_READY|" + file_path, "").strip()
            if not full_reply:
                full_reply = "🧠 已为您生成文档核心知识关系图谱。"
        except Exception as e:
            print(f"解析KGRAPH失败：{e}")

    # 3. 解析总结报告
    if "SUMMARY_READY|" in tool_output_str or "SUMMARY_READY|" in full_reply:
        try:
            source_str = tool_output_str if "SUMMARY_READY|" in tool_output_str else full_reply
            marker = source_str.index("SUMMARY_READY|")
            download_path = source_str[marker + 14:].strip().split('\n')[0].strip()
            full_reply = full_reply.replace("SUMMARY_READY|" + download_path, "").strip() + "\n\n📝 **总结文件已生成！**"
        except Exception as e:
            print(f"解析SUMMARY失败：{e}")

    # 4. 解析音乐预览
    if "MUSIC_PREVIEW|" in tool_output_str or "MUSIC_PREVIEW|" in full_reply:
        try:
            source_str = tool_output_str if "MUSIC_PREVIEW|" in tool_output_str else full_reply
            # 使用正则提取完整的 MUSIC_PREVIEW 标记
            match = re.search(r"MUSIC_PREVIEW\|([^|]+)\|([^|]+)\|(.+)", source_str)
            if match:
                name = match.group(1)
                artist = match.group(2)
                preview_url = match.group(3).strip()
                # 验证文件是否存在
                if os.path.exists(preview_url):
                    download_path = preview_url
                    full_reply = full_reply.replace(f"MUSIC_PREVIEW|{name}|{artist}|{preview_url}", "").strip() + f"\n\n🎵 **音乐试听：** {name} - {artist}"
                else:
                    preview_url = None
        except Exception as e:
            print(f"解析音乐预览失败：{e}")

    history[-1]["content"] = full_reply
    tasks_dict[current_name] = history

    # 保存当前任务的图表、音频和下载状态
    if current_name not in task_visuals:
        task_visuals[current_name] = {"heatmap": None, "audio": None, "download": None}
    task_visuals[current_name]["heatmap"] = heatmap_html
    task_visuals[current_name]["audio"] = preview_url
    task_visuals[current_name]["download"] = download_path

    # 使用 gr.update 同时设置组件值和可见性
    audio_output = gr.update(value=preview_url, visible=(preview_url is not None))
    download_output = gr.update(value=download_path, visible=(download_path is not None))

    yield "", list(history), update_radio(tasks_dict,
                                          current_name), current_name, audio_output, heatmap_html, download_output


def new_chat(tasks_dict, task_visuals):
    """新建对话."""
    tasks_dict["新对话"] = []
    task_visuals["新对话"] = {"heatmap": None, "audio": None, "download": None}
    audio_output = gr.update(value=None, visible=False)
    download_output = gr.update(value=None, visible=False)
    return "", [], gr.update(choices=list(tasks_dict.keys()), value="新对话"), "新对话", audio_output, None, download_output


def switch_chat(name, tasks_dict, task_visuals):
    """切换对话."""
    if isinstance(name, list):
        name = name[0] if name else ""
    if name in tasks_dict:
        visuals = task_visuals.get(name, {"heatmap": None, "audio": None, "download": None})
        heatmap = visuals.get("heatmap")
        audio = visuals.get("audio")
        download = visuals.get("download")
        
        # 验证文件路径是否存在
        if audio and not os.path.exists(audio):
            audio = None
        if download and not os.path.exists(download):
            download = None
            
        audio_output = gr.update(value=audio, visible=audio is not None)
        download_output = gr.update(value=download, visible=download is not None)
        return list(tasks_dict[name]), gr.update(choices=list(tasks_dict.keys()), value=name), name, audio_output, heatmap, download_output
    audio_output = gr.update(value=None, visible=False)
    download_output = gr.update(value=None, visible=False)
    return [], gr.update(choices=list(tasks_dict.keys())), "", audio_output, None, download_output


def delete_chat(tasks_dict, current_name, task_visuals):
    """删除对话."""
    if current_name and current_name in tasks_dict:
        del tasks_dict[current_name]
        if current_name in task_visuals:
            del task_visuals[current_name]
    keys = list(tasks_dict.keys())
    new_name = keys[-1] if keys else ""
    new_history = list(tasks_dict[new_name]) if new_name and new_name in tasks_dict else []
    if new_name and new_name in task_visuals:
        visuals = task_visuals[new_name]
        heatmap = visuals.get("heatmap")
        audio = visuals.get("audio")
        download = visuals.get("download")
    else:
        heatmap = None
        audio = None
        download = None
    audio_output = gr.update(value=audio, visible=audio is not None)
    download_output = gr.update(value=download, visible=download is not None)
    return new_history, gr.update(choices=keys, value=new_name), new_name, audio_output, heatmap, download_output


def clear_chat(tasks_dict, current_name, task_visuals):
    """清空当前对话."""
    if current_name and current_name in tasks_dict:
        tasks_dict[current_name] = []
        if current_name in task_visuals:
            task_visuals[current_name] = {"heatmap": None, "audio": None, "download": None}
    audio_output = gr.update(value=None, visible=False)
    download_output = gr.update(value=None, visible=False)
    return [], "", audio_output, None, download_output


# ==================== Gradio界面 ====================

with gr.Blocks() as demo:
    tasks = gr.State({"新对话": []})
    cur_name = gr.State("新对话")
    # 每个任务对应的图表和下载文件状态
    task_visuals = gr.State({})

    with gr.Row():
        # 左侧任务列表
        with gr.Column(scale=0, min_width=180):
            gr.Markdown("### 对话列表")
            new_btn = gr.Button("+ 新对话", variant="primary")
            task_list = gr.Radio(choices=["新对话"], value="新对话", label="")
            del_btn = gr.Button("删除任务", variant="stop")
            clr_btn = gr.Button("清空记录")

        # 右侧主区域
        with gr.Column(scale=1):
            gr.Markdown("# AI助手\n天气查询 · 数据计算 · 邮件发送 · 股票分析 · 文件分析 · 知识图谱 · 搜题库")

            # 展示区
            with gr.Row():
                audio_player = gr.Audio(label="🎵 音乐试听", visible=False)
                download_file = gr.File(label="📥 下载区", interactive=False, visible=False)

            with gr.Row():
                heatmap_display = gr.HTML(label="📊 图表展示区")

            chatbot = gr.Chatbot(height=500)

            with gr.Row():
                file_upload = gr.File(label="📁 上传文件", type="filepath", scale=1)
                txt = gr.Textbox(placeholder="请输入指令...", show_label=False, scale=4)
                send = gr.Button("发送", variant="primary", scale=1)

    # 事件绑定
    new_btn.click(new_chat, [tasks, task_visuals], [txt, chatbot, task_list, cur_name, audio_player, heatmap_display, download_file])
    task_list.select(switch_chat, [task_list, tasks, task_visuals], [chatbot, task_list, cur_name, audio_player, heatmap_display, download_file])

    # 发送消息（带文件上传）
    send.click(chat_stream, [txt, chatbot, tasks, cur_name, file_upload, task_visuals],
               [txt, chatbot, task_list, cur_name, audio_player, heatmap_display, download_file])
    txt.submit(chat_stream, [txt, chatbot, tasks, cur_name, file_upload, task_visuals],
               [txt, chatbot, task_list, cur_name, audio_player, heatmap_display, download_file])


    # 文件上传自动处理
    def on_file_uploaded(file_path, current_chatbot, tasks_dict, current_name):
        if not file_path:
            audio_output = gr.update(value=None, visible=False)
            download_output = gr.update(value=None, visible=False)
            return gr.update(), current_chatbot, gr.update(), current_name, audio_output, None, download_output
        auto_prompt = "我已经上传了文件，请先帮我详细总结这份文档讲了什么，然后再为它生成知识图谱。"
        audio_output = gr.update(value=None, visible=False)
        download_output = gr.update(value=None, visible=False)
        return gr.update(value=auto_prompt), current_chatbot, gr.update(), current_name, audio_output, None, download_output


    file_upload.upload(on_file_uploaded, [file_upload, chatbot, tasks, cur_name],
                       [txt, chatbot, task_list, cur_name, audio_player, heatmap_display, download_file])

    del_btn.click(delete_chat, [tasks, cur_name, task_visuals], [chatbot, task_list, cur_name, audio_player, heatmap_display, download_file])
    clr_btn.click(clear_chat, [tasks, cur_name, task_visuals], [chatbot, txt, audio_player, heatmap_display, download_file])

demo.launch(allowed_paths=[tempfile.gettempdir()])
